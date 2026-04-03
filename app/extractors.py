from __future__ import annotations

import io
import json
from pathlib import Path
from typing import Any, BinaryIO

from .models import ContentBlock, NormalizedDocument
from .normalizer import join_block_texts, markdown_table, normalize_text


class UnsupportedFileTypeError(ValueError):
    pass


def detect_file_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    return ext or "txt"


def _get_handler(file_type: str):
    handlers = {
        "pdf": extract_pdf,
        "docx": extract_docx,
        "md": extract_text,
        "txt": extract_text,
        "json": extract_json,
        "html": extract_html,
        "htm": extract_html,
        "csv": extract_table,
        "xlsx": extract_table,
        "xls": extract_table,
        "pptx": extract_pptx,
    }
    handler = handlers.get(file_type)
    if not handler:
        raise UnsupportedFileTypeError(f"Unsupported file type: {file_type}")
    return handler


def extract(path: Path) -> NormalizedDocument:
    file_type = detect_file_type(path)
    handler = _get_handler(file_type)
    with path.open("rb") as stream:
        return handler(stream, path.name, file_type, str(path))


def extract_bytes(filename: str, data: bytes, source_path: str | None = None) -> NormalizedDocument:
    path = Path(filename)
    file_type = detect_file_type(path)
    handler = _get_handler(file_type)
    stream = io.BytesIO(data)
    return handler(stream, path.name, file_type, source_path)


def extract_text(stream: BinaryIO, source_name: str, file_type: str, source_path: str | None = None) -> NormalizedDocument:
    text = stream.read().decode("utf-8", errors="ignore")
    block = ContentBlock(type="text", text=text, order=0)
    return _build_doc(source_name, file_type, [block], source_path=source_path)


def extract_pdf(stream: BinaryIO, source_name: str, file_type: str, source_path: str | None = None) -> NormalizedDocument:
    from pypdf import PdfReader

    reader = PdfReader(stream)
    blocks: list[ContentBlock] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        blocks.append(ContentBlock(type="page", text=text, order=idx - 1, metadata={"page": idx}))
    metadata: dict[str, Any] = {}
    if reader.metadata:
        metadata.update({k.lstrip("/"): v for k, v in reader.metadata.items()})
    return _build_doc(source_name, file_type, blocks, metadata, source_path=source_path)


def extract_docx(stream: BinaryIO, source_name: str, file_type: str, source_path: str | None = None) -> NormalizedDocument:
    from docx import Document

    doc = Document(stream)
    blocks: list[ContentBlock] = []
    order = 0
    for para in doc.paragraphs:
        text = para.text or ""
        if not text.strip():
            continue
        block_type = "heading" if para.style and para.style.name and para.style.name.startswith("Heading") else "paragraph"
        blocks.append(ContentBlock(type=block_type, text=text, order=order))
        order += 1
    for table in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_text = markdown_table(rows)
        if table_text:
            blocks.append(ContentBlock(type="table", text=table_text, order=order))
            order += 1
    return _build_doc(source_name, file_type, blocks, source_path=source_path)


def extract_json(stream: BinaryIO, source_name: str, file_type: str, source_path: str | None = None) -> NormalizedDocument:
    raw = stream.read().decode("utf-8", errors="ignore")
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        data = raw
    pretty = json.dumps(data, ensure_ascii=False, indent=2) if not isinstance(data, str) else data
    block = ContentBlock(type="json", text=pretty, order=0)
    return _build_doc(source_name, file_type, [block], source_path=source_path)


def extract_html(stream: BinaryIO, source_name: str, file_type: str, source_path: str | None = None) -> NormalizedDocument:
    from bs4 import BeautifulSoup

    raw = stream.read().decode("utf-8", errors="ignore")
    soup = BeautifulSoup(raw, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text("\n")
    block = ContentBlock(type="text", text=text, order=0)
    return _build_doc(source_name, file_type, [block], source_path=source_path)


def extract_table(stream: BinaryIO, source_name: str, file_type: str, source_path: str | None = None) -> NormalizedDocument:
    import pandas as pd

    blocks: list[ContentBlock] = []
    order = 0
    if file_type == "csv":
        df = pd.read_csv(stream)
        blocks.append(ContentBlock(type="table", text=df.to_markdown(index=False), order=order))
    else:
        xls = pd.ExcelFile(stream)
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            blocks.append(
                ContentBlock(
                    type="sheet",
                    text=df.to_markdown(index=False),
                    order=order,
                    metadata={"sheet": sheet_name},
                )
            )
            order += 1
    return _build_doc(source_name, file_type, blocks, source_path=source_path)


def extract_pptx(stream: BinaryIO, source_name: str, file_type: str, source_path: str | None = None) -> NormalizedDocument:
    from pptx import Presentation

    prs = Presentation(stream)
    blocks: list[ContentBlock] = []
    order = 0
    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slide_text = "\n".join(texts)
        blocks.append(ContentBlock(type="slide", text=slide_text, order=order, metadata={"slide": slide_idx}))
        order += 1
    return _build_doc(source_name, file_type, blocks, source_path=source_path)


def _build_doc(
    source_name: str,
    file_type: str,
    blocks: list[ContentBlock],
    metadata: dict[str, Any] | None = None,
    source_path: str | None = None,
) -> NormalizedDocument:
    normalized_blocks = []
    for block in blocks:
        normalized_blocks.append(
            ContentBlock(
                type=block.type,
                text=normalize_text(block.text),
                order=block.order,
                metadata=block.metadata,
            )
        )
    text = join_block_texts(normalized_blocks)
    return NormalizedDocument(
        source_name=source_name,
        source_path=source_path,
        file_type=file_type,
        title=Path(source_name).stem,
        text=text,
        content_blocks=normalized_blocks,
        metadata=metadata or {},
    )
